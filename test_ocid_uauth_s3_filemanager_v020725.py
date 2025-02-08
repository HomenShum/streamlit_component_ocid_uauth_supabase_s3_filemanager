import streamlit as st
import supabase
import boto3
import os
from botocore.exceptions import NoCredentialsError, ClientError
import math # For pagination
import pandas as pd
from io import BytesIO
import base64
from pptx import Presentation # Ensure pptx is installed: pip install python-pptx


# Initialize Supabase client
supabase_url = st.secrets['supabase']['SUPABASE_URL']
supabase_key = st.secrets['supabase']['SUPABASE_KEY']
supabase_client = supabase.create_client(supabase_url, supabase_key)

# Supabase Storage details (from secrets or config)
SUPABASE_S3_BUCKET_NAME = st.secrets['supabase']['SUPABASE_S3_BUCKET_NAME']

# Initialize boto3 client for S3 (Supabase Storage)
s3_client = boto3.client(
    's3',
    endpoint_url=st.secrets['supabase']['SUPABASE_S3_ENDPOINT_URL'],
    region_name=st.secrets['supabase']['SUPABASE_S3_BUCKET_REGION'],
    aws_access_key_id=st.secrets['supabase']['SUPABASE_S3_BUCKET_ACCESS_KEY'],
    aws_secret_access_key=st.secrets['supabase']['SUPABASE_S3_BUCKET_SECRET_KEY']
)

KEY_PREFIX = "s3_file_manager" # To avoid session state conflicts
ITEMS_PER_PAGE_OPTIONS = [5, 10, 25, 50, 100] # Pagination options

# --- Session State Initialization ---
def _init_session_state():
    user_root_path = f"{st.experimental_user.name}" if st.experimental_user.is_logged_in else "" # No trailing slash here
    user_root_path = user_root_path + "/" if user_root_path else "" # Add it back if not empty
    if KEY_PREFIX + '_current_path' not in st.session_state:
        st.session_state[KEY_PREFIX + '_current_path'] = user_root_path

    if KEY_PREFIX + '_previous_path' not in st.session_state:
        st.session_state[KEY_PREFIX + '_previous_path'] = "" # Or None initially
    if KEY_PREFIX + '_show_new_folder_input' not in st.session_state:
        st.session_state[KEY_PREFIX + '_show_new_folder_input'] = False
    if KEY_PREFIX + '_show_rename_folder_input' not in st.session_state: # Not implementing rename in this version for simplicity
        st.session_state[KEY_PREFIX + '_show_rename_folder_input'] = False
    if KEY_PREFIX + '_show_upload' not in st.session_state:
        st.session_state[KEY_PREFIX + '_show_upload'] = False # Initially hide upload, show on button click
    if KEY_PREFIX + '_selected_folders' not in st.session_state:
        st.session_state[KEY_PREFIX + '_selected_folders'] = []
    if KEY_PREFIX + '_selected_files' not in st.session_state:
        st.session_state[KEY_PREFIX + '_selected_files'] = []
    if KEY_PREFIX + '_selected_files_in_folders' not in st.session_state: # NEW: For files in selected folders
        st.session_state[KEY_PREFIX + '_selected_files_in_folders'] = []
    if KEY_PREFIX + '_new_folder_name' not in st.session_state:
        st.session_state[KEY_PREFIX + '_new_folder_name'] = ""
    if KEY_PREFIX + '_rename_folder_name' not in st.session_state: # Not implementing rename in this version for simplicity
        st.session_state[KEY_PREFIX + '_rename_folder_name'] = ""
    if KEY_PREFIX + '_upload_success' not in st.session_state:
        st.session_state[KEY_PREFIX + '_upload_success'] = []
    if KEY_PREFIX + '_upload_progress' not in st.session_state:
        st.session_state[KEY_PREFIX + '_upload_progress'] = 0
    if KEY_PREFIX + '_current_page' not in st.session_state: # For pagination
        st.session_state[KEY_PREFIX + '_current_page'] = 1
    if KEY_PREFIX + '_items_per_page' not in st.session_state: # For pagination
        st.session_state[KEY_PREFIX + '_items_per_page'] = ITEMS_PER_PAGE_OPTIONS[1] # Default to 10 items per page
    if KEY_PREFIX + '_delete_confirmation' not in st.session_state:
        st.session_state[KEY_PREFIX + '_delete_confirmation'] = {} # Dict to hold confirmation state for each item

# @st.cache_data(show_spinner=False, ttl=10)
def list_files_in_folder(folder_path):
    """Lists files within a given S3 folder path."""
    folders, files = list_s3_files(prefix=folder_path) # Use existing list_s3_files
    full_file_paths = files  # Files from list_s3_files already have the correct path
    return full_file_paths

# @st.cache_data(show_spinner=False, ttl=10) # Caching for performance - adjust ttl as needed
def list_s3_files(prefix=""):
    """Lists files and folders in an S3 bucket under a given prefix."""
    try:
        response = s3_client.list_objects_v2(Bucket=SUPABASE_S3_BUCKET_NAME, Prefix=prefix, Delimiter='/') # Delimiter for folders
        files = []
        folders = []
        if 'CommonPrefixes' in response: # Folders are returned in CommonPrefixes
            for prefix_info in response['CommonPrefixes']:
                folders.append(prefix_info['Prefix'])
        if 'Contents' in response: # Files are in Contents
            for obj in response['Contents']:
                if not obj['Key'].endswith('/'): # Exclude folder "placeholders"
                    files.append(obj['Key'])
        return folders, files

    except NoCredentialsError:
        st.error("AWS credentials not available.")
        return [], []
    except ClientError as e:
        st.error(f"Error accessing S3: {e}")
        return [], []

def upload_file_to_s3(file, s3_key):
    """Uploads a file-like object to S3."""
    try:
        s3_client.upload_fileobj(file, SUPABASE_S3_BUCKET_NAME, s3_key)
        return True
    except NoCredentialsError:
        st.error("AWS credentials not available.")
        return False
    except ClientError as e:
        st.error(f"Error uploading to S3: {e}")
        return False

def download_file_from_s3(s3_key):
    """Downloads a file from S3 and returns its content as bytes."""
    try:
        file_obj = s3_client.get_object(Bucket=SUPABASE_S3_BUCKET_NAME, Key=s3_key)
        file_content = file_obj['Body'].read()
        # st.success(f"File content retrieved from S3 for: {s3_key}") # No need for success here, handled in UI
        return file_content
    except NoCredentialsError:
        st.error("AWS credentials not available.")
        return None
    except ClientError as e:
        st.error(f"Error downloading file from S3: {e}")
        return None

def delete_file_from_s3(s3_key):
    """Deletes a file from S3 after checking if it exists."""
    sanitized_key = sanitize_path(s3_key)  # Sanitize the S3 key
    try:
        # Check if the object exists
        head_response = s3_client.head_object(Bucket=SUPABASE_S3_BUCKET_NAME, Key=sanitized_key)
        print(f"Head response: {head_response}")
        # Proceed to delete
        delete_response = s3_client.delete_object(Bucket=SUPABASE_S3_BUCKET_NAME, Key=sanitized_key)
        print(f"Delete response: {delete_response}")
        return True, s3_key # Return True and the s3_key of the deleted file
    except ClientError as e:
        if e.response['Error']['Code'] == '404':
            st.error(f"File not found: {sanitized_key}")
        else:
            st.error(f"Error deleting from S3: {e}")
        return False, None # Return False and None if deletion fails

def create_s3_folder(s3_folder_key):
    """Creates an empty folder (object with '/' suffix) in S3."""
    sanitized_folder_key = sanitize_path(s3_folder_key)  # Sanitize the folder key
    try:
        s3_client.put_object(Bucket=SUPABASE_S3_BUCKET_NAME, Key=f"{sanitized_folder_key}/") # Keys for folders must end with '/'
        return True
    except NoCredentialsError:
        st.error("AWS credentials not available.")
        return False
    except ClientError as e:
        st.error(f"Error creating S3 folder: {e}")
        print(f"Error details: {e}") # Print error details for debugging
        return False

def sanitize_path(path):
    return path.strip('/').replace('//', '/')

def delete_s3_folder(s3_folder_prefix):
    """Recursively deletes a folder and all its contents from S3."""
    sanitized_prefix = sanitize_path(s3_folder_prefix)
    print(f"Sanitized prefix for deletion: {sanitized_prefix}")

    paginator = s3_client.get_paginator('list_objects_v2')
    pages = paginator.paginate(Bucket=SUPABASE_S3_BUCKET_NAME, Prefix=sanitized_prefix)

    for page in pages:
        if 'Contents' in page:
            for obj in page['Contents']:
                try:
                    s3_client.delete_object(Bucket=SUPABASE_S3_BUCKET_NAME, Key=obj['Key'])
                    print(f"Deleted object: {obj['Key']}")
                except ClientError as e:
                    print(f"Error deleting object {obj['Key']}: {e}")

    # Attempt to delete the folder "placeholder" object
    placeholder_key = sanitized_prefix + '/'
    try:
        s3_client.delete_object(Bucket=SUPABASE_S3_BUCKET_NAME, Key=placeholder_key)
        print(f"Deleted placeholder object: {placeholder_key}")
    except ClientError as e:
        if e.response['Error']['Code'] == 'NoSuchKey':
            print(f"Placeholder object '{placeholder_key}' not found (which is okay).")
        else:
            print(f"Error deleting placeholder object '{placeholder_key}': {e}")

    return True



def _format_size(size: int) -> str:
    """Format file size in human-readable format."""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size < 1024:
            return f"{size:.1f} {unit}"
        size /= 1024
    return f"{size:.1f} TB"

def _render_pagination(total_items: int):
    """Render pagination controls."""
    total_pages = math.ceil(total_items / st.session_state[KEY_PREFIX + '_items_per_page'])

    if total_pages > 1:
        col1, col2, col3, col4, col5 = st.columns([1, 1, 2, 1, 1])
        current_page = st.session_state[KEY_PREFIX + '_current_page']

        with col1:
            if st.button("⏮️", disabled=current_page == 1, key=f"{KEY_PREFIX}first"):
                st.session_state[KEY_PREFIX + '_current_page'] = 1
                st.rerun()

        with col2:
            if st.button("◀️", disabled=current_page == 1, key=f"{KEY_PREFIX}prev"):
                st.session_state[KEY_PREFIX + '_current_page'] -= 1
                st.rerun()
        with col3:
            page_options = list(range(1, total_pages + 1))
            selected_page = st.selectbox(
                "Go to page",
                options=page_options,
                index=current_page - 1,
                key=f"{KEY_PREFIX}page_select",
                label_visibility="collapsed"
            )
            if selected_page != current_page:
                st.session_state[KEY_PREFIX + '_current_page'] = selected_page
                st.rerun()

        with col4:
            if st.button("▶️", disabled=current_page == total_pages, key=f"{KEY_PREFIX}next"):
                st.session_state[KEY_PREFIX + '_current_page'] += 1
                st.rerun()

        with col5:
            if st.button("⏭️", disabled=current_page == total_pages, key=f"{KEY_PREFIX}last"):
                st.session_state[KEY_PREFIX + '_current_page'] = total_pages
                st.rerun()

def get_file_type_from_extension(filename: str) -> str:
    """Extracts file type from filename extension."""
    _, ext = os.path.splitext(filename)
    if ext:
        return ext[1:].upper()  # Remove the dot and uppercase
    return "Unknown"


def render_folder_management_ui():
    current_path = st.session_state[KEY_PREFIX + '_current_path']
    previous_path = st.session_state[KEY_PREFIX + '_previous_path']
    root_path = f"{st.experimental_user.name}/" if st.experimental_user.is_logged_in else ""

    print(f"Current Path in render_folder_management_ui: {current_path}") # Debug print

    st.markdown(f"**S3 Bucket:** `{SUPABASE_S3_BUCKET_NAME}`")
    st.markdown(f"**Root Folder (Your Files):** `{root_path if root_path else 'root of bucket'}`") # Clarify root

    path_components = [comp for comp in current_path.split('/') if comp] # Split and remove empty strings
    full_path = root_path if root_path else "" # Start building path from root

    path_display_cols_count = len(path_components) + (1 if current_path != root_path and current_path != "" else 0)
    path_display_cols = st.columns(max(1, path_display_cols_count))  # Ensure at least 1 column

    # "Back" button if not at root
    if current_path != root_path and current_path != "":
        if st.button("🔙 Back to Previous Folder", key="back_button"):
            st.session_state[KEY_PREFIX + '_current_path'] = st.session_state[KEY_PREFIX + '_previous_path'] # Go back to previous path

            # Simplify previous path update: Get parent directory of the *current* previous_path
            current_prev_path = st.session_state[KEY_PREFIX + '_previous_path']
            if current_prev_path: # Only update if there's a previous path
                st.session_state[KEY_PREFIX + '_previous_path'] = os.path.dirname(current_prev_path.rstrip('/')) + "/" # Get parent and ensure trailing slash
            else:
                st.session_state[KEY_PREFIX + '_previous_path'] = "" # If no previous path, set to root

            st.rerun()

    # Ensure we do not exceed the number of columns
    for i, component in enumerate(path_components):
        full_path = os.path.join(full_path, component)  # Reconstruct path
        if i < len(path_display_cols):  # Check if index is within range
            with path_display_cols[i]:  # Use i directly for the current column
                if i < len(path_components) - 1:  # Make path components clickable except the last one
                    if st.button(component, key=f"path_comp_btn_{i}", help=f"Go to '{full_path}'"):
                        st.session_state[KEY_PREFIX + '_previous_path'] = st.session_state[KEY_PREFIX + '_current_path']  # Update previous path
                        st.session_state[KEY_PREFIX + '_current_path'] = full_path + "/"  # Ensure trailing slash for folder prefix
                        st.rerun()
                else:
                    st.markdown(f"**{component}**")  # Current folder as bold text
        else:
            st.warning(f"Path display column index out of range for component: {component}. This should not happen, please report.")  # Debugging warning

    folders_in_folder, files_in_folder = list_s3_files(prefix=current_path)
    items = sorted([{'name': os.path.basename(f_name), 'path': f_name, 'is_directory': False, 'size': None} for f_name in files_in_folder] + # USE f_name directly for file path
                   [{'name': os.path.basename(f_prefix.rstrip('/')), 'path': f_prefix, 'is_directory': True, 'size': None} for f_prefix in folders_in_folder],
                   key=lambda x: (not x['is_directory'], x['name'].lower()))

    start_idx = (st.session_state[KEY_PREFIX + '_current_page'] - 1) * st.session_state[KEY_PREFIX + '_items_per_page']
    end_idx = start_idx + st.session_state[KEY_PREFIX + '_items_per_page']
    paginated_items = items[start_idx:end_idx]

    if paginated_items:
        for item in paginated_items:
            col_sel, col_name, col_size, col_actions = st.columns([0.5, 4, 2, 3])
            with col_sel:
                if item['is_directory']:
                    folder_selected = item['path'] in st.session_state[KEY_PREFIX + '_selected_folders']
                    checkbox_selected = st.checkbox("Select Folder", key=f"folder_checkbox_{item['path']}", value=folder_selected, label_visibility="collapsed")
                    if checkbox_selected:
                        if item['path'] not in st.session_state[KEY_PREFIX + '_selected_folders']:
                            st.session_state[KEY_PREFIX + '_selected_folders'].append(item['path'])
                            # NEW: Add files in selected folder to selected_files_in_folders and deselect from _selected_files
                            files_in_sel_folder = list_files_in_folder(item['path'])
                            st.session_state[KEY_PREFIX + '_selected_files_in_folders'].extend(files_in_sel_folder)
                            st.session_state[KEY_PREFIX + '_selected_files'] = [
                                f for f in st.session_state[KEY_PREFIX + '_selected_files']
                                if not f.startswith(item['path']) # Remove files in selected folder from _selected_files
                            ]
                    else:
                        if item['path'] in st.session_state[KEY_PREFIX + '_selected_folders']:
                            st.session_state[KEY_PREFIX + '_selected_folders'].remove(item['path'])
                            # NEW: Remove files of deselected folder from selected_files_in_folders and _selected_files
                            files_in_sel_folder = list_files_in_folder(item['path'])
                            st.session_state[KEY_PREFIX + '_selected_files_in_folders'] = [f for f in st.session_state[KEY_PREFIX + '_selected_files_in_folders'] if f not in files_in_sel_folder]
                            st.session_state[KEY_PREFIX + '_selected_files'] = [
                                f for f in st.session_state[KEY_PREFIX + '_selected_files']
                                if not f.startswith(item['path']) # Remove files in deselected folder from _selected_files
                            ]


                else:
                    # --- Modified file_selected condition ---
                    file_selected = (item['path'] in st.session_state[KEY_PREFIX + '_selected_files'] or item['path'] in st.session_state[KEY_PREFIX + '_selected_files_in_folders'])
                    if st.checkbox("Select File", key=f"file_checkbox_{item['path']}", value=file_selected, label_visibility="collapsed"):
                        if item['path'] not in st.session_state[KEY_PREFIX + '_selected_files']:
                            st.session_state[KEY_PREFIX + '_selected_files'].append(item['path'])
                            # --- Auto-select folder on file select logic ---
                            current_folder_path = os.path.dirname(item['path']) + "/" # Normalize folder path with trailing slash
                            files_in_current_folder_view_all = [i for i in paginated_items if not i['is_directory'] and i['path'].startswith(current_folder_path)]
                            total_files_in_view_in_current_folder = len(files_in_current_folder_view_all)
                            selected_files_in_view_in_current_folder_check = 0
                            for file_item_check in files_in_current_folder_view_all:
                                if file_item_check['path'] in st.session_state[KEY_PREFIX + '_selected_files']:
                                    selected_files_in_view_in_current_folder_check += 1
                            if total_files_in_view_in_current_folder > 0 and total_files_in_view_in_current_folder == selected_files_in_view_in_current_folder_check:
                                folder_item_to_auto_select = next((fd for fd in paginated_items if fd['is_directory'] and fd['path'] == current_folder_path), None)
                                if folder_item_to_auto_select and folder_item_to_auto_select['path'] not in st.session_state[KEY_PREFIX + '_selected_folders']:
                                    st.session_state[KEY_PREFIX + '_selected_folders'].append(folder_item_to_auto_select['path'])

                    else:
                        if item['path'] in st.session_state[KEY_PREFIX + '_selected_files']:
                            st.session_state[KEY_PREFIX + '_selected_files'].remove(item['path'])
                            # --- Auto-deselect folder on file deselect logic ---
                            current_folder_path = os.path.dirname(item['path']) + "/" # Normalize folder path with trailing slash
                            if current_folder_path in st.session_state[KEY_PREFIX + '_selected_folders']:
                                st.session_state[KEY_PREFIX + '_selected_folders'].remove(current_folder_path) # Deselect folder if a file is deselected
                            # --- End auto-deselect folder on file deselect logic ---


            with col_name:
                if item['is_directory']:
                    if st.button(f"📁 {item['name']}", key=f"open_folder_btn_{item['path']}", use_container_width=True, help=f"Open Folder: {item['name']}"):
                        st.session_state[KEY_PREFIX + '_previous_path'] = current_path
                        st.session_state[KEY_PREFIX + '_current_path'] = os.path.join(current_path, item['name']) + "/" # Just join and *then* add trailing slash
                        st.session_state[KEY_PREFIX + '_current_page'] = 1
                        st.rerun()
                else:
                    st.markdown(f"📄 {item['name']}")

            with col_size:
                if not item['is_directory']:
                    # Get file size (inefficient to get size for each file in list, consider optimizing if needed)
                    try:
                        response = s3_client.head_object(Bucket=SUPABASE_S3_BUCKET_NAME, Key=item['path'])
                        file_size = response['ContentLength']
                        st.text(_format_size(file_size))
                    except ClientError as e:
                        st.text("Size N/A") # Handle errors getting file size
                else:
                    st.empty() # No size for folders
            with col_actions:
                if not item['is_directory']:
                    if st.button("Download ⬇️", key=f"download_file_btn_{item['path']}", use_container_width=True, help=f"Download File: {item['name']}"):
                        file_content = download_file_from_s3(item['path'])
                        if file_content:
                            st.download_button(
                                label="Click to Download",
                                data=file_content,
                                file_name=item['name'],
                                mime="application/octet-stream",
                                key=f"download_button_{item['path']}"
                            )
                            st.success(f"File download ready: {item['name']}", icon="⬇️")
                        else:
                            st.error("Failed to download file content.")
                    if st.button("Delete 🗑️", key=f"delete_btn_{item['path']}", use_container_width=True, help=f"Delete {'Folder' if item['is_directory'] else 'File'}: {item['name']}"):
                        if item['is_directory']:
                            if delete_s3_folder(item['path']):
                                st.success(f"Folder '{item['name']}' deleted.")
                                # --- Folder Delete Update ---
                                folder_prefix_to_delete = item['path']
                                # Remove deleted folder from selected folders
                                if folder_prefix_to_delete in st.session_state[KEY_PREFIX + '_selected_folders']:
                                    st.session_state[KEY_PREFIX + '_selected_folders'].remove(folder_prefix_to_delete)
                                # Filter out files that were in the deleted folder from selected_files_in_folders
                                st.session_state[KEY_PREFIX + '_selected_files_in_folders'] = [
                                    f for f in st.session_state[KEY_PREFIX + '_selected_files_in_folders']
                                    if not f.startswith(folder_prefix_to_delete)
                                ]
                                # Filter out files that were in the deleted folder from selected_files (standalone selected files)
                                st.session_state[KEY_PREFIX + '_selected_files'] = [
                                    f for f in st.session_state[KEY_PREFIX + '_selected_files']
                                    if not f.startswith(folder_prefix_to_delete)
                                ]
                                # --- End Folder Delete Update ---
                                st.rerun()
                            else:
                                st.error(f"Failed to delete folder '{item['name']}'.")
                        else:
                            deleted_successfully, deleted_key = delete_file_from_s3(item['path']) # Capture deleted_key
                            if deleted_successfully:
                                st.success(f"File '{item['name']}' deleted.")
                                # --- START OF FILE DELETE UPDATE ---
                                if deleted_key in st.session_state[KEY_PREFIX + '_selected_files']:
                                    st.session_state[KEY_PREFIX + '_selected_files'].remove(deleted_key)

                                # Refresh selected_files_in_folders if the deleted file was in a selected folder
                                selected_folder_containing_deleted_file = next((folder for folder in st.session_state[KEY_PREFIX + '_selected_folders'] if deleted_key.startswith(folder)), None)
                                if selected_folder_containing_deleted_file:
                                    files_in_sel_folder = list_files_in_folder(selected_folder_containing_deleted_file)
                                    st.session_state[KEY_PREFIX + '_selected_files_in_folders'] = files_in_sel_folder # Directly replace to refresh list
                                # --- END OF FILE DELETE UPDATE ---
                                st.rerun()
                            else:
                                st.error(f"Failed to delete file '{item['name']}'.")

    elif folders_in_folder or files_in_folder: # If folder is not empty but no items to display on current page
        st.info(f"No items to display on page {st.session_state[KEY_PREFIX + '_current_page']}. Please use pagination controls to navigate.")
    else:
        st.info("This folder is empty.")
    _render_pagination(len(items))  # Pagination below the list

    # Display Selected Paths Section in DataFrame
    st.subheader("Selected Items:")
    if st.session_state[KEY_PREFIX + '_selected_folders'] or st.session_state[KEY_PREFIX + '_selected_files'] or st.session_state[KEY_PREFIX + '_selected_files_in_folders']:
        selected_folders = st.session_state[KEY_PREFIX + '_selected_folders']
        selected_files = st.session_state[KEY_PREFIX + '_selected_files']
        selected_files_in_folders = st.session_state[KEY_PREFIX + '_selected_files_in_folders']

        data = []
        for folder in selected_folders:
            folder_name = os.path.basename(folder.rstrip('/')) # Get folder name without trailing slash
            data.append({"Folder": folder_name, "Type": "Folder", "File Name": folder_name, "Path": folder})
            files_in_folder = [f for f in selected_files_in_folders if f.startswith(folder)]
            for file in files_in_folder:
                file_name = os.path.basename(file)
                file_type = get_file_type_from_extension(file_name)
                data.append({"Folder": folder_name, "Type": file_type, "File Name": file_name, "Path": file})

        # Handle standalone selected files (not in selected folders)
        standalone_selected_files = [
            f for f in selected_files
            if not any(f.startswith(folder) for folder in selected_folders)
        ]
        for file in standalone_selected_files:
            file_folder_path = os.path.dirname(file) # Get the folder path
            file_folder_name = os.path.basename(file_folder_path) if file_folder_path else "N/A" # Extract folder name or N/A if root
            if file_folder_path == "" or file_folder_name == st.experimental_user.name or file_folder_name == "": # Handle root or user root edge cases
                file_folder_name = "N/A"

            file_name = os.path.basename(file)
            file_type = get_file_type_from_extension(file_name)
            data.append({"Folder": file_folder_name, "Type": file_type, "File Name": file_name, "Path": file}) # Use correct file_folder_name

        df = pd.DataFrame(data)
        if not df.empty: # Check if DataFrame is not empty before displaying
            st.dataframe(df[["Folder", "Type", "File Name", "Path"]], use_container_width=True, hide_index=True) # Order columns and hide index
        else:
            st.info("No items to display in DataFrame (this should not happen if selected items exist).") # Debugging info
    else:
        st.info("No folders or files selected.")


def render_action_buttons():
    col1, col2, col3 = st.columns([2, 2, 3]) # Adjust column widths as needed

    with col1:
        if st.button("➕ New Folder"):
            st.session_state[KEY_PREFIX + '_show_new_folder_input'] = True

        if st.session_state[KEY_PREFIX + '_show_new_folder_input']:
            new_folder_name = st.text_input("Folder name:", key=KEY_PREFIX + '_new_folder_name_input')
            create_button = st.button("✅ Create", key="create_folder_btn")
            cancel_new_folder_button = st.button("❌ Cancel", key="cancel_new_folder_btn")

            if create_button:
                if new_folder_name:
                    new_s3_folder_key = os.path.join(st.session_state[KEY_PREFIX + '_current_path'], new_folder_name)
                    new_s3_folder_key = new_s3_folder_key + "/" # Ensure trailing slash for folder key
                    if create_s3_folder(new_s3_folder_key):
                        st.success(f"Folder '{new_folder_name}' created in '{st.session_state[KEY_PREFIX + '_current_path']}'!")
                        st.session_state[KEY_PREFIX + '_show_new_folder_input'] = False # Hide input after creation
                        st.session_state[KEY_PREFIX + '_new_folder_name'] = "" # Clear input
                        st.rerun()
                    else:
                        st.error(f"Failed to create folder '{new_folder_name}'. Check logs.")
                else:
                    st.warning("Please enter a folder name.")
            if cancel_new_folder_button:
                st.session_state[KEY_PREFIX + '_show_new_folder_input'] = False
                st.session_state[KEY_PREFIX + '_new_folder_name'] = "" # Clear input

    with col2:
        if st.button("📤 Upload Files"): # Changed button text to icon + text
            st.session_state[KEY_PREFIX + '_show_upload'] = not st.session_state[KEY_PREFIX + '_show_upload'] # Toggle upload section

    with col3:
        if st.button("🗑️ Delete Folders"):
            if st.session_state[KEY_PREFIX + '_selected_folders']:
                folders_to_delete = st.session_state[KEY_PREFIX + '_selected_folders']
                folders_to_delete = st.session_state[KEY_PREFIX + '_selected_folders']
                for folder_prefix in folders_to_delete:
                    if delete_s3_folder(folder_prefix):
                        st.success(f"Folder '{os.path.basename(folder_prefix.rstrip('/'))}' deleted.")
                        # --- Folder Delete Update ---
                        folder_prefix_to_delete = folder_prefix
                        # Remove deleted folder from selected folders
                        if folder_prefix_to_delete in st.session_state[KEY_PREFIX + '_selected_folders']:
                            st.session_state[KEY_PREFIX + '_selected_folders'].remove(folder_prefix_to_delete)
                        # Filter out files that were in the deleted folder from selected_files_in_folders
                        st.session_state[KEY_PREFIX + '_selected_files_in_folders'] = [
                            f for f in st.session_state[KEY_PREFIX + '_selected_files_in_folders']
                            if not f.startswith(folder_prefix_to_delete)
                        ]
                        # Filter out files that were in the deleted folder from selected_files (standalone selected files)
                        st.session_state[KEY_PREFIX + '_selected_files'] = [
                            f for f in st.session_state[KEY_PREFIX + '_selected_files']
                            if not f.startswith(folder_prefix_to_delete)
                        ]
                        # --- End Folder Delete Update ---
                    else:
                        st.error(f"Failed to delete folder '{os.path.basename(folder_prefix.rstrip('/'))}'. Check logs.")
                st.session_state[KEY_PREFIX + '_selected_folders'] = []
                st.session_state[KEY_PREFIX + '_selected_files_in_folders'] = []
                st.rerun()
            else:
                st.warning("No folders selected for deletion.")


def render_upload_section():
    if st.session_state[KEY_PREFIX + '_show_upload']:
        with st.expander("📤 Upload Files", expanded=True): # Expander for upload section
            uploaded_files = st.file_uploader("Choose files to upload", accept_multiple_files=True, key=KEY_PREFIX + "_file_uploader")
            if uploaded_files:
                for uploaded_file in uploaded_files:
                    s3_key_upload = os.path.join(st.session_state[KEY_PREFIX + '_current_path'], uploaded_file.name) # Construct S3 key with folder path
                    if upload_file_to_s3(uploaded_file, s3_key_upload):
                        st.success(f"File '{uploaded_file.name}' uploaded to '{s3_key_upload}'")
                        if st.session_state[KEY_PREFIX + '_current_path'] in st.session_state[KEY_PREFIX + '_selected_folders']:
                            # Re-list files in the current (selected) folder and update selected_files_in_folders
                            files_in_current_folder = list_files_in_folder(st.session_state[KEY_PREFIX + '_current_path'])
                            # Remove files from current folder and then add the new list to avoid duplicates
                            st.session_state[KEY_PREFIX + '_selected_files_in_folders'] = [
                                f for f in st.session_state[KEY_PREFIX + '_selected_files_in_folders']
                                if not f.startswith(st.session_state[KEY_PREFIX + '_current_path'])
                            ]
                            st.session_state[KEY_PREFIX + '_selected_files_in_folders'].extend(files_in_current_folder)
                    else:
                        st.error(f"Failed to upload '{uploaded_file.name}'")
                st.session_state[KEY_PREFIX + '_show_upload'] = False # Hide upload section after upload
                st.rerun() # Refresh file list after upload

def render_items_per_page_selector():
    items_per_page = st.selectbox(
        "Items per page",
        options=ITEMS_PER_PAGE_OPTIONS,
        key=KEY_PREFIX + "items_per_page_selector",
        label_visibility="collapsed" # Hide label to save space
    )
    if items_per_page != st.session_state[KEY_PREFIX + '_items_per_page']:
        st.session_state[KEY_PREFIX + '_items_per_page'] = items_per_page
        st.session_state[KEY_PREFIX + '_current_page'] = 1 # Reset to page 1 when items per page changes
        st.rerun()

@st.fragment
def sidebar_content_fragment_st_file_manager_component():
    st.title("Lite S3 File Manager (Supabase Storage)")

    # Ensure user is logged in
    if not st.experimental_user.is_logged_in:
        st.button("Log in with Google", on_click=st.login)
        st.stop()

    _init_session_state() # Initialize session state

    st.header("Browse & Manage S3 Files") # More specific header

    with st.container(border=True): # Container for action buttons and file listing
        col_header, col_pagination_selector = st.columns([5, 2]) # Adjust ratio as needed for header and selector
        with col_header:
            st.subheader("File & Folder Actions") # Moved subheader into container
        with col_pagination_selector:
            render_items_per_page_selector() # Items per page selector in line with header

        render_action_buttons()
        render_upload_section() # Render upload section below actions
        render_folder_management_ui() # File/folder listing


def main():
    with st.sidebar:
        sidebar_content_fragment_st_file_manager_component()

    selected_documents = st.session_state[KEY_PREFIX + '_selected_files'] + st.session_state[KEY_PREFIX + '_selected_files_in_folders']

    if selected_documents:
        st.subheader("Selected Documents Preview:")
        # No changes needed in the preview section to accommodate multiple files.
        # The existing tab logic already handles displaying multiple files in separate tabs.
        tab_names = [f"{f[:10]}...{f[-10:]}" if len(f) > 25 else f for f in map(os.path.basename, selected_documents)]
        tabs = st.tabs(tab_names)

        for idx, file_path in enumerate(selected_documents):
            with tabs[idx]:
                if len(selected_documents) > 1:
                    with st.popover("Tips 💡"):
                        st.write("Scroll horizontally to view and select all documents.")

                corpus_path = os.path.dirname(file_path)
                corpus_name = os.path.basename(corpus_path)
                st.write(f"**Corpus:** {corpus_name}, **Document:** {os.path.basename(file_path)}") # Display corpus name

                if file_path.endswith(".pdf"):
                    st.write(f"File type: PDF")
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        base64_pdf = base64.b64encode(file_content).decode('utf-8')
                        pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800px" type="application/pdf"></iframe>'
                        st.markdown(pdf_display, unsafe_allow_html=True)
                    else:
                        st.error("Failed to load PDF content.")

                elif file_path.endswith((".csv", ".tsv")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            df = pd.read_csv(BytesIO(file_content))
                            st.dataframe(df)
                        except Exception as e:
                            st.error(f"Error reading CSV/TSV: {e}")
                    else:
                        st.error("Failed to load CSV/TSV content.")
                elif file_path.endswith((".xlsx", ".xls")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            if file_path.endswith(".xlsx"):
                                df = pd.read_excel(BytesIO(file_content), engine='openpyxl') # Specify engine for xlsx
                            elif file_path.endswith(".xls"):
                                df = pd.read_excel(BytesIO(file_content), engine='xlrd') # Specify engine for xls
                            else: # Fallback if somehow extension is not recognized
                                df = pd.read_excel(BytesIO(file_content)) # Let pandas try to infer
                            st.dataframe(df)
                        except Exception as e:
                            st.error(f"Error reading Excel file: {e}")
                    else:
                        st.error("Failed to load Excel content.")
                elif file_path.endswith((".doc", ".docx", ".txt", ".html", ".md", ".rtf")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            text_content = file_content.decode('utf-8', errors='ignore') # Handle encoding issues
                            with st.container(border=True):
                                st.markdown(text_content)
                        except Exception as e:
                            st.error(f"Error displaying text-based file: {e}")
                    else:
                        st.error("Failed to load text-based file content.")
                elif file_path.endswith((".json", ".xml")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            df = pd.read_json(BytesIO(file_content)) if file_path.endswith(".json") else pd.read_xml(BytesIO(file_content))
                            st.dataframe(df)
                        except Exception as e:
                            st.error(f"Error reading JSON/XML: {e}")
                    else:
                        st.error("Failed to load JSON/XML content.")
                elif file_path.endswith((".mp4", ".avi", ".mov", ".webm", ".mkv")):
                    # Need to create a temporary local file for st.video to work with S3 content efficiently
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            temp_video_file = BytesIO(file_content)
                            st.video(temp_video_file)
                        except Exception as e:
                            st.error(f"Error displaying video: {e}")
                    else:
                        st.error("Failed to load video content.")
                elif file_path.endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            st.image(BytesIO(file_content))
                        except Exception as e:
                            st.error(f"Error displaying image: {e}")
                    else:
                        st.error("Failed to load image content.")
                elif file_path.endswith((".mp3", ".wav", ".ogg", ".flac")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            st.audio(BytesIO(file_content))
                        except Exception as e:
                            st.error(f"Error displaying audio: {e}")
                    else:
                        st.error("Failed to load audio content.")
                elif file_path.endswith((".py", ".js", ".java", ".cpp", ".cs", ".rb")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            code_content = file_content.decode('utf-8', errors='ignore')
                            with st.container(border=True):
                                st.code(code_content, language=file_path.split('.')[-1])
                        except Exception as e:
                            st.error(f"Error displaying code file: {e}")
                    else:
                        st.error("Failed to load code file content.")
                elif file_path.endswith((".ppt", ".pptx")):
                    file_content = download_file_from_s3(file_path)
                    if file_content:
                        try:
                            prs = Presentation(BytesIO(file_content))
                            pdf_buffer = BytesIO()
                            c = canvas.Canvas(pdf_buffer, pagesizes=letter)

                            for slide in prs.slides:
                                c.drawString(100, 750, f"Slide {prs.slides.index(slide) + 1}")
                                for shape in slide.shapes:
                                    if hasattr(shape, 'text'):
                                        c.drawString(100, 700, shape.text[:50])  # Truncate long texts
                                c.showPage()

                            c.save()
                            pdf_bytes = base64.b64encode(pdf_buffer.getvalue()).decode('utf-8') # Corrected base64 encoding
                            pdf_display = f'<iframe src="data:application/pdf;base64,{pdf_bytes}" width="100%" height="500px" type="application/pdf"></iframe>'
                            st.markdown(pdf_display, unsafe_allow_html=True)
                        except ImportError:
                            st.write("PowerPoint file detected. Preview not available due to missing dependencies.")
                        except Exception as e:
                            st.error(f"Error displaying PowerPoint: {e}")
                    else:
                        st.error("Failed to load PowerPoint content.")

                elif file_path.endswith(".zip"):
                    st.write("ZIP file detected. Contents cannot be displayed directly.")
                elif file_path.endswith((".accdb", ".mdb")):
                    st.write("Access database file detected. Preview not available.")
                elif file_path.endswith(".mpp"):
                    st.write("Microsoft Project file detected. Preview not available.")
                elif file_path.endswith((".one", ".onetoc2")):
                    st.write("OneNote file detected. Preview not available.")
                elif file_path.endswith(".vsd"):
                    st.write("Visio drawing file detected. Preview not available.")
                else:
                    st.write("Unsupported document type")


if __name__ == "__main__":
    main()
